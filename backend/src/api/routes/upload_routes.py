"""Upload routes — PDF/DOCX/TXT ingestion with full pipeline."""
from fastapi import APIRouter, Depends, UploadFile, File, Request
from core.security import get_current_user
from core.limiter import limiter
from core.constants import RATE_LIMIT_UPLOAD
from utils.validators import validate_file
from rag.ingestion.pdf_loader import load_pdf
from rag.ingestion.docx_loader import load_docx
from rag.ingestion.text_loader import load_text
from rag.ingestion.embedding_pipeline import run_ingestion_pipeline
from rag.vectorstores.pinecone_store import get_vector_store
from database.mongodb.repositories.document_repository import save_document_metadata
from core.logger import get_logger

router = APIRouter()
logger = get_logger(__name__)


@router.post("/upload")
@limiter.limit(RATE_LIMIT_UPLOAD)
async def upload_document(
    request: Request,
    file: UploadFile = File(...),
    current_user=Depends(get_current_user),
):
    user_id = str(current_user["_id"])
    filename = file.filename

    validate_file(filename, file.size or 0)

    content = await file.read()

    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    IMAGE_EXTS = {"png", "jpg", "jpeg", "gif", "webp", "bmp", "tiff", "svg"}
    SHEET_EXTS = {"xlsx", "xls", "csv"}
    PRES_EXTS = {"pptx", "ppt"}
    TEXT_EXTS = {"txt", "md", "mdx", "rst", "html", "htm", "json", "jsonl",
                 "xml", "yaml", "yml", "csv",
                 "py", "js", "ts", "jsx", "tsx", "java", "cpp", "c", "cs",
                 "go", "rs", "rb", "php", "sh", "sql", "rtf", "odt", "doc"}

    if ext == "pdf":
        docs = await load_pdf(content, filename)
        # Check if PDF is scanned (empty or very short text)
        total_text = "".join(d.page_content for d in docs).strip()
        if len(total_text) < 100:
            logger.info(f"PDF '{filename}' appears to be scanned (extracted text length: {len(total_text)}). Running vision extraction...")
            from agents.langgraph.nodes.vision_extraction_node import extract_scanned_pdf
            docs = await extract_scanned_pdf(content, filename)
    elif ext == "docx":
        docs = await load_docx(content, filename)
    elif ext in IMAGE_EXTS:
        from agents.langgraph.nodes.vision_extraction_node import extract_vision_data
        from langchain_core.documents import Document
        import base64
        
        extracted_content = await extract_vision_data(
            content=content,
            filename=filename,
            file_type=ext,
            extraction_type="auto"
        )
        
        docs = [Document(
            page_content=extracted_content,
            metadata={"filename": filename, "source": filename, "file_type": ext, "page": 0,
                      "data_uri": "data:image/" + ext + ";base64," + base64.b64encode(content[:4096]).decode()},
        )]
    elif ext in SHEET_EXTS and ext != "csv":
        try:
            import io
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                rows.append(f"=== Sheet: {ws.title} ===")
                for row in ws.iter_rows(values_only=True):
                    rows.append("\t".join(str(c) if c is not None else "" for c in row))
            docs = await load_text("\n".join(rows).encode(), filename)
        except Exception as xlsx_err:
            logger.warning(f"Excel parse failed ({xlsx_err}), falling back to text loader")
            docs = await load_text(content, filename)
    elif ext in PRES_EXTS:
        try:
            import io
            from pptx import Presentation as PptxPresentation
            prs = PptxPresentation(io.BytesIO(content))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
                slides_text.append(f"=== Slide {i} ===\n" + "\n".join(texts))
            docs = await load_text("\n".join(slides_text).encode(), filename)
        except Exception as pptx_err:
            logger.warning(f"PPTX parse failed ({pptx_err}), falling back to text loader")
            docs = await load_text(content, filename)
    else:
        docs = await load_text(content, filename)

    vector_store = get_vector_store()
    chunk_count = await run_ingestion_pipeline(
        documents=docs,
        user_id=user_id,
        filename=filename,
        vector_store=vector_store,
        use_semantic_chunking=True,
    )

    await save_document_metadata(
        user_id=user_id,
        filename=filename,
        file_type=ext,
        chunk_count=chunk_count,
        file_bytes=content,
    )

    try:
        from services.workflow_scheduler import trigger_event_workflows
        asyncio.create_task(
            trigger_event_workflows(
                event_type="document_uploaded",
                user_id=user_id,
                payload={"filename": filename, "file_type": ext, "chunks": chunk_count},
            )
        )
    except Exception as e:
        logger.warning(f"Failed to trigger event workflow for upload: {e}")

    return {
        "message": "Document ingested and indexed successfully",
        "filename": filename,
        "chunks": chunk_count,
    }


from fastapi.responses import StreamingResponse
import json
import asyncio

@router.post("/upload/stream")
@limiter.limit(RATE_LIMIT_UPLOAD)
async def upload_document_stream(
    request: Request,
    file: UploadFile = File(...),
    current_user=Depends(get_current_user),
):
    user_id = str(current_user["_id"])
    filename = file.filename
    content = await file.read()

    async def event_generator():
        try:
            yield f"data: {json.dumps({'stage': 'parsing', 'progress': 10})}\n\n"
            validate_file(filename, len(content))
            ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

            IMAGE_EXTS = {"png", "jpg", "jpeg", "gif", "webp", "bmp", "tiff", "svg"}
            SHEET_EXTS = {"xlsx", "xls"}
            PRES_EXTS = {"pptx", "ppt"}

            if ext == "pdf":
                docs = await load_pdf(content, filename)
                # Check if PDF is scanned (empty or very short text)
                total_text = "".join(d.page_content for d in docs).strip()
                if len(total_text) < 100:
                    logger.info(f"PDF '{filename}' appears to be scanned (extracted text length: {len(total_text)}). Running vision extraction...")
                    from agents.langgraph.nodes.vision_extraction_node import extract_scanned_pdf
                    docs = await extract_scanned_pdf(content, filename)
            elif ext == "docx":
                docs = await load_docx(content, filename)
            elif ext in IMAGE_EXTS:
                from agents.langgraph.nodes.vision_extraction_node import extract_vision_data
                from langchain_core.documents import Document
                
                extracted_content = await extract_vision_data(
                    content=content,
                    filename=filename,
                    file_type=ext,
                    extraction_type="auto"
                )
                
                docs = [Document(
                    page_content=extracted_content,
                    metadata={"filename": filename, "source": filename, "file_type": ext, "page": 0},
                )]
            elif ext in SHEET_EXTS:
                try:
                    import io
                    import openpyxl
                    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
                    rows = []
                    for ws in wb.worksheets:
                        rows.append(f"=== Sheet: {ws.title} ===")
                        for row in ws.iter_rows(values_only=True):
                            rows.append("\t".join(str(c) if c is not None else "" for c in row))
                    docs = await load_text("\n".join(rows).encode(), filename)
                except Exception as xlsx_err:
                    logger.warning(f"Excel parse failed ({xlsx_err}), falling back to text loader")
                    docs = await load_text(content, filename)
            elif ext in PRES_EXTS:
                try:
                    import io
                    from pptx import Presentation as PptxPresentation
                    prs = PptxPresentation(io.BytesIO(content))
                    slides_text = []
                    for i, slide in enumerate(prs.slides, 1):
                        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
                        slides_text.append(f"=== Slide {i} ===\n" + "\n".join(texts))
                    docs = await load_text("\n".join(slides_text).encode(), filename)
                except Exception as pptx_err:
                    logger.warning(f"PPTX parse failed ({pptx_err}), falling back to text loader")
                    docs = await load_text(content, filename)
            else:
                docs = await load_text(content, filename)
            
            yield f"data: {json.dumps({'stage': 'chunking', 'progress': 30})}\n\n"
            from rag.ingestion.embedding_pipeline import clean_documents, semantic_chunk, enrich_metadata
            from rag.vectorstores.pinecone_store import get_vector_store

            vector_store = get_vector_store()
            
            docs = clean_documents(docs)
            
            chunks = semantic_chunk(docs)
            yield f"data: {json.dumps({'stage': 'embedding', 'progress': 60, 'chunks': len(chunks)})}\n\n"
            
            chunks = enrich_metadata(chunks, user_id=user_id, filename=filename)
            
            yield f"data: {json.dumps({'stage': 'indexing', 'progress': 80})}\n\n"
            await asyncio.to_thread(vector_store.add_documents, chunks)
            chunk_count = len(chunks)

            await save_document_metadata(
                user_id=user_id,
                filename=filename,
                file_type=ext,
                chunk_count=chunk_count,
                file_bytes=content,
            )

            yield f"data: {json.dumps({'stage': 'done', 'progress': 100, 'chunks': chunk_count})}\n\n"

        except Exception as e:
            logger.error(f"Streaming upload failed: {e}")
            yield f"data: {json.dumps({'error': str(e)})}\n\n"

    return StreamingResponse(event_generator(), media_type="text/event-stream")


@router.get("")
async def list_documents(current_user=Depends(get_current_user)):
    from database.mongodb.repositories.document_repository import get_user_documents
    user_id = str(current_user["_id"])
    return await get_user_documents(user_id)


@router.get("/{document_id}/download")
async def download_document(document_id: str, current_user=Depends(get_current_user)):
    from bson.objectid import ObjectId
    from fastapi.responses import Response
    from fastapi import HTTPException
    from database.mongodb.connection import get_database
    from core.constants import COLLECTION_DOCUMENTS

    user_id = str(current_user["_id"])
    db = get_database()
    
    doc = await db[COLLECTION_DOCUMENTS].find_one({
        "_id": ObjectId(document_id),
        "user_id": user_id
    })
    
    if not doc or "file_bytes" not in doc:
        raise HTTPException(status_code=404, detail="Document file not found")
        
    filename = doc.get("filename", "document")
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    
    media_types = {
        "pdf": "application/pdf",
        "png": "image/png",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "gif": "image/gif",
        "webp": "image/webp",
        "svg": "image/svg+xml",
        "txt": "text/plain",
        "md": "text/plain",
        "csv": "text/csv",
        "json": "application/json",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    }
    media_type = media_types.get(ext, "application/octet-stream")
    
    return Response(
        content=doc["file_bytes"],
        media_type=media_type,
        headers={
            "Content-Disposition": f'inline; filename="{filename}"'
        }
    )


@router.delete("/{document_id}")
async def delete_document(document_id: str, current_user=Depends(get_current_user)):
    from bson.objectid import ObjectId
    from fastapi import HTTPException
    from database.mongodb.connection import get_database
    from core.constants import COLLECTION_DOCUMENTS
    from database.mongodb.repositories.document_repository import delete_document_metadata
    from rag.vectorstores.pinecone_store import delete_by_filename
    
    user_id = str(current_user["_id"])
    db = get_database()
    
    doc = await db[COLLECTION_DOCUMENTS].find_one({
        "_id": ObjectId(document_id),
        "user_id": user_id
    })
    
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
        
    filename = doc.get("filename")
    
    db_success = await delete_document_metadata(user_id, document_id)
    
    pinecone_success = False
    if filename:
        try:
            pinecone_success = delete_by_filename(user_id, filename)
        except Exception as pine_err:
            logger.warning(f"Could not delete {filename} from Pinecone: {pine_err}")
    
    if db_success or pinecone_success:
        return {"message": "Document deleted"}
        
    raise HTTPException(status_code=404, detail="Document could not be deleted")
